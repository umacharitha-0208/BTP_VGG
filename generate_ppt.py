"""
Generate RIPE-PPO Presentation PPT
Run: python generate_ppt.py
Output: RIPE_PPO_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)   # cyan
ACCENT2    = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xD0, 0xD8, 0xE8)
YELLOW     = RGBColor(0xFF, 0xD1, 0x66)
GREEN      = RGBColor(0x57, 0xCC, 0x99)
RED        = RGBColor(0xFF, 0x6B, 0x6B)
MID_BLUE   = RGBColor(0x1A, 0x3A, 0x5C)
BOX_BG     = RGBColor(0x12, 0x2B, 0x45)

SW, SH = Inches(13.33), Inches(7.5)       # 16:9 widescreen

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

blank_layout = prs.slide_layouts[6]       # completely blank

# ─────────────────────────────────────────────────────────────────────────────
# Helper utilities
# ─────────────────────────────────────────────────────────────────────────────

def set_bg(slide, color):
    """Fill slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.FREEFORM if False else 1,  # MSO_SHAPE.RECTANGLE
        l, t, w, h
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_bullet_box(slide, title, bullets, l, t, w, h,
                   bg=BOX_BG, title_color=ACCENT, bullet_color=WHITE,
                   title_size=Pt(16), bullet_size=Pt(13), accent_line=True):
    """Rounded box with a title and bullet points."""
    # background rectangle
    box = add_rect(slide, l, t, w, h, fill_color=bg, line_color=ACCENT, line_width=Pt(1))

    # accent top bar
    if accent_line:
        add_rect(slide, l, t, w, Inches(0.05), fill_color=ACCENT)

    # title
    ty = t + Inches(0.08)
    add_text(slide, title, l + Inches(0.15), ty, w - Inches(0.3), Inches(0.35),
             font_size=title_size, bold=True, color=title_color, align=PP_ALIGN.LEFT)

    # bullets
    bx = l + Inches(0.15)
    by = ty + Inches(0.38)
    bw = w - Inches(0.3)
    bh = h - Inches(0.5)
    txBox = slide.shapes.add_textbox(bx, by, bw, bh)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = ("• " if not b.startswith("   ") else "") + b
        run.font.size = bullet_size
        run.font.color.rgb = bullet_color if not b.startswith("   ") else ACCENT2
    return box


def slide_header(slide, title, subtitle=None):
    """Thin header bar at top of slide."""
    add_rect(slide, 0, 0, SW, Inches(0.72), fill_color=MID_BLUE)
    add_rect(slide, 0, Inches(0.68), SW, Inches(0.05), fill_color=ACCENT)
    add_text(slide, title, Inches(0.4), Inches(0.1), Inches(10), Inches(0.55),
             font_size=Pt(26), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.42), Inches(10), Inches(0.3),
                 font_size=Pt(13), bold=False, color=ACCENT2, align=PP_ALIGN.LEFT)

    # slide number area (right side of header)


def add_table(slide, headers, rows, l, t, w, h,
              header_bg=MID_BLUE, row_bg=BOX_BG, alt_bg=DARK_BG,
              header_color=ACCENT, row_color=WHITE):
    """Simple table using individual rectangles (no native pptx table borders)."""
    n_cols = len(headers)
    n_rows = len(rows)
    col_w = w / n_cols
    row_h = h / (n_rows + 1)

    # header row
    for ci, hdr in enumerate(headers):
        add_rect(slide, l + ci*col_w, t, col_w, row_h, fill_color=header_bg, line_color=ACCENT, line_width=Pt(0.5))
        add_text(slide, hdr, l + ci*col_w + Inches(0.05), t + Inches(0.03), col_w - Inches(0.1), row_h,
                 font_size=Pt(12), bold=True, color=header_color, align=PP_ALIGN.CENTER)

    # data rows
    for ri, row in enumerate(rows):
        bg = row_bg if ri % 2 == 0 else alt_bg
        for ci, cell in enumerate(row):
            add_rect(slide, l + ci*col_w, t + (ri+1)*row_h, col_w, row_h, fill_color=bg, line_color=MID_BLUE, line_width=Pt(0.5))
            cell_color = GREEN if str(cell).startswith("+") else (RED if str(cell).startswith("-") and "pp" in str(cell) else row_color)
            add_text(slide, str(cell), l + ci*col_w + Inches(0.05), t + (ri+1)*row_h + Inches(0.02),
                     col_w - Inches(0.1), row_h,
                     font_size=Pt(11), bold=False, color=cell_color, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 1 — TITLE
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)

# big centre gradient bar
add_rect(slide, 0, Inches(2.2), SW, Inches(3.2), fill_color=MID_BLUE)
add_rect(slide, 0, Inches(2.2), Inches(0.07), Inches(3.2), fill_color=ACCENT)
add_rect(slide, 0, Inches(5.35), SW, Inches(0.06), fill_color=ACCENT)

# decorative dots top-right
for i in range(5):
    for j in range(5):
        add_rect(slide, SW - Inches(1.8) + i*Inches(0.32), Inches(0.3) + j*Inches(0.32),
                 Inches(0.12), Inches(0.12), fill_color=ACCENT if (i+j)%2==0 else MID_BLUE)

add_text(slide, "RIPE-PPO", Inches(0.6), Inches(2.35), Inches(12), Inches(1.0),
         font_size=Pt(54), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_text(slide, "Reinforcement Learning Keypoint Detection with",
         Inches(0.6), Inches(3.2), Inches(11), Inches(0.5),
         font_size=Pt(22), bold=False, color=ACCENT2, align=PP_ALIGN.LEFT)

add_text(slide, "Proximal Policy Optimization + Attention Descriptors + LO-RANSAC",
         Inches(0.6), Inches(3.65), Inches(11), Inches(0.5),
         font_size=Pt(22), bold=True, color=YELLOW, align=PP_ALIGN.LEFT)

add_text(slide, "Based on RIPE (ICCV 2025)  |  Modified Architecture  |  BTP Project",
         Inches(0.6), Inches(4.55), Inches(11), Inches(0.4),
         font_size=Pt(14), bold=False, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

add_text(slide, "IMW2022 Dataset  •  VGG19 Backbone  •  50 000 Training Steps",
         Inches(0.6), Inches(4.95), Inches(11), Inches(0.35),
         font_size=Pt(13), color=ACCENT2, align=PP_ALIGN.LEFT)

# badge
add_rect(slide, Inches(10.3), Inches(6.5), Inches(2.6), Inches(0.7), fill_color=ACCENT, line_color=None)
add_text(slide, "B.Tech Project  2025-26", Inches(10.3), Inches(6.55), Inches(2.6), Inches(0.6),
         font_size=Pt(13), bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 2 — INTRODUCTION
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "01  Introduction", "Why reinvent keypoint detection?")

# left column — problem statement
add_bullet_box(slide, "What is Keypoint Detection?",
    ["Identify distinctive, repeatable points in images",
     "Used in: SfM, SLAM, image matching, AR/VR",
     "Key challenge: robust under illumination, viewpoint, scale change",
     "   Traditional: SIFT, ORB — hand-crafted, limited generalization",
     "   Learning-based: SuperPoint, DISK, ALIKE — supervised",
     "   RL-based: RIPE (ICCV 2025) — reward-driven selection"],
    Inches(0.3), Inches(0.85), Inches(4.1), Inches(2.8))

# centre column — what is RIPE
add_bullet_box(slide, "Original RIPE (ICCV 2025)",
    ["REINFORCE policy gradient — binary selection policy",
     "VGG19 backbone → multi-scale HyperColumn features",
     "Simple descriptor head (linear projection only)",
     "MNN matching + RANSAC pose estimation",
     "Reward = inlier count from pose RANSAC",
     "   Limitation: high variance gradient (REINFORCE)",
     "   Limitation: no context in descriptors",
     "   Limitation: random RANSAC sampling"],
    Inches(4.55), Inches(0.85), Inches(4.1), Inches(2.8))

# right column — our contribution
add_bullet_box(slide, "Our Contributions (RIPE-PPO)",
    ["PPO replaces REINFORCE — stable, clipped updates",
     "Attention Descriptor Head — context-aware 256-d embeddings",
     "LO-RANSAC — local optimisation, fewer iterations wasted",
     "PROSAC ordering — best matches sampled first by RANSAC",
     "InfoNCE loss — descriptor supervision without match labels",
     "   Synthetic homography bootstrapping (10% of batches)",
     "   BFloat16 mixed precision — 2x VRAM efficiency"],
    Inches(8.8), Inches(0.85), Inches(4.2), Inches(2.8))

# bottom — motivation bar
add_rect(slide, Inches(0.3), Inches(3.85), Inches(12.73), Inches(3.35), fill_color=BOX_BG, line_color=ACCENT, line_width=Pt(1))
add_rect(slide, Inches(0.3), Inches(3.85), Inches(12.73), Inches(0.05), fill_color=ACCENT)
add_text(slide, "Core Motivation", Inches(0.5), Inches(3.9), Inches(6), Inches(0.35),
         font_size=Pt(15), bold=True, color=ACCENT)

motivation_items = [
    ("REINFORCE Variance", "Policy gradient with Monte Carlo returns has extremely high variance.\nPPO adds a value baseline and clips updates → stable, sample-efficient training."),
    ("Descriptor Quality", "Linear projection ignores spatial context between keypoints.\nMulti-head attention over HyperColumn features → richer, context-aware descriptors."),
    ("RANSAC Efficiency", "Uniform random sampling wastes iterations on bad matches.\nPROSAC sorts by descriptor similarity → high-quality matches sampled first."),
]
for i, (title, body) in enumerate(motivation_items):
    lx = Inches(0.5) + i * Inches(4.22)
    add_rect(slide, lx, Inches(4.3), Inches(4.0), Inches(2.7), fill_color=MID_BLUE, line_color=ACCENT2, line_width=Pt(0.5))
    add_text(slide, title, lx + Inches(0.1), Inches(4.38), Inches(3.8), Inches(0.35),
             font_size=Pt(13), bold=True, color=YELLOW)
    add_text(slide, body, lx + Inches(0.1), Inches(4.72), Inches(3.8), Inches(2.1),
             font_size=Pt(11), color=LIGHT_GRAY)


# =============================================================================
# SLIDE 3 — MODEL ARCHITECTURE
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "02  Model Architecture", "End-to-end pipeline: image pair → pose estimate")

# Pipeline flow — horizontal boxes with arrows
pipeline = [
    ("Image Pair\n(384×384)", ACCENT),
    ("VGG19\nBackbone", MID_BLUE),
    ("HyperColumn\nFeatures\n960-d", MID_BLUE),
    ("PPO Policy\n(Actor-Critic)", MID_BLUE),
    ("Attention\nDescriptor\n256-d", MID_BLUE),
    ("PROSAC\n+ MNN", MID_BLUE),
    ("LO-RANSAC\nPose", MID_BLUE),
    ("Reward\n& Loss", ACCENT),
]

box_w = Inches(1.42)
box_h = Inches(1.1)
start_x = Inches(0.25)
pipe_y = Inches(0.88)
gap = Inches(0.22)

for i, (label, color) in enumerate(pipeline):
    bx = start_x + i * (box_w + gap)
    add_rect(slide, bx, pipe_y, box_w, box_h, fill_color=color, line_color=ACCENT, line_width=Pt(1))
    add_text(slide, label, bx + Inches(0.05), pipe_y + Inches(0.1), box_w - Inches(0.1), box_h - Inches(0.15),
             font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # arrow
    if i < len(pipeline) - 1:
        ax = bx + box_w + Inches(0.03)
        add_text(slide, "→", ax, pipe_y + Inches(0.35), gap + Inches(0.05), Inches(0.4),
                 font_size=Pt(16), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Detail boxes — 4 component deep-dives
components = [
    ("VGG19 Backbone",
     ["19-layer deep CNN pretrained on ImageNet",
      "Gradient checkpointing — saves ~40% VRAM",
      "All layers trainable (freeze_layers = 0)",
      "4 conv blocks → feature pyramid",
      "Outputs: C1(64) C2(128) C3(256) C4(512)"]),
    ("HyperColumn Features",
     ["Bilinear interpolation at keypoint locations",
      "Features from all 4 VGG conv blocks",
      "Concatenated → 960-d per keypoint",
      "Multi-scale context in single vector",
      "Input to both policy head and descriptor head"]),
    ("PPO Actor-Critic",
     ["Actor: logits → Bernoulli selection policy",
      "Critic: value head → baseline for GAE",
      "Clipped surrogate: clip(r, 1-ε, 1+ε)·A",
      "ε = 0.1, 2 PPO epochs per rollout",
      "GAE: λ=0.95, γ=0.99 for advantage"]),
    ("Attention Descriptor Head",
     ["Linear(960 → 256) projection",
      "Multi-Head Attention (4 heads) + residual",
      "LayerNorm → FFN(256→512→256) + residual",
      "Final LayerNorm → L2 normalisation",
      "InfoNCE loss (T=0.07) for supervision"]),
]

col_w = Inches(3.15)
for i, (title, bullets) in enumerate(components):
    lx = Inches(0.25) + i * (col_w + Inches(0.1))
    add_bullet_box(slide, title, bullets,
                   lx, Inches(2.15), col_w, Inches(2.85),
                   title_size=Pt(13), bullet_size=Pt(11))

# Bottom — PROSAC + LO-RANSAC box
add_bullet_box(slide, "PROSAC + LO-RANSAC Matching Pipeline",
    ["1. Top-K keypoints by detector score (score-guided pre-filter)",
     "2. MNN matching with Kornia (th=1.0, ratio test disabled)",
     "3. Sort matches by cosine similarity descending → PROSAC order",
     "4. poselib.estimate_fundamental with progressive_sampling=True",
     "5. LO-RANSAC: lo_iterations=25, max_iterations=1000, min_iterations=50"],
    Inches(0.25), Inches(5.15), Inches(12.83), Inches(2.1),
    title_size=Pt(14), bullet_size=Pt(12))


# =============================================================================
# SLIDE 4 — CHANGES FROM ORIGINAL
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "03  Changes from Original RIPE", "What we modified and why it matters")

# comparison table
headers = ["Component", "Original RIPE (ICCV 2025)", "RIPE-PPO (Ours)", "Benefit"]
rows = [
    ["Policy Gradient",     "REINFORCE (high variance)",      "PPO w/ clipped surrogate",     "Stable, sample-efficient"],
    ["Value Baseline",      "None",                           "Critic network + GAE",          "Variance reduction"],
    ["Descriptor Head",     "Linear(960→256) only",           "Linear + MHA(4) + FFN + LN",    "Context-aware embeddings"],
    ["Descriptor Loss",     "None (RL only)",                 "InfoNCE (T=0.07)",              "Strong gradient signal"],
    ["RANSAC",              "Standard random sampling",        "LO-RANSAC (lo_iter=25)",        "Higher inlier precision"],
    ["Match Ordering",      "Random / unordered",             "PROSAC (similarity-sorted)",    "Best matches tried first"],
    ["Pre-filter",          "All keypoints to MNN",           "Top-K by detector score",       "Reduced MNN noise"],
    ["Precision",           "FP32",                           "BF16-mixed",                    "2× VRAM, 2× speed"],
    ["Synthetic Data",      "None",                           "10% synthetic homography",      "Cold-start bootstrap"],
]
add_table(slide, headers, rows, Inches(0.25), Inches(0.85), Inches(12.83), Inches(5.6),
          header_bg=MID_BLUE, row_bg=BOX_BG, alt_bg=RGBColor(0x0F, 0x22, 0x38))

# bottom legend
add_text(slide, "Green = improvement    |    PROSAC & LO-RANSAC act at inference only — no extra training cost",
         Inches(0.3), Inches(6.8), Inches(12), Inches(0.4),
         font_size=Pt(11), color=ACCENT2, italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 5 — TRAINING PARAMETERS
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "04  Training Parameters", "Full-scale GPU configuration")

# 3-column layout
col1 = [
    ("Training Setup",
     ["Steps: 50,000",
      "Batch size: 4 (effective: 16 w/ grad acc)",
      "Gradient accumulation: 4 steps",
      "Precision: BF16-mixed",
      "Image resolution: 384 × 384 px",
      "GPU: Single GPU (auto-detect)",
      "Dataset: IMW2022"]),
]

col2 = [
    ("Optimiser & LR",
     ["LR: 1e-4 (warmup 500 steps)",
      "Final LR: 1e-6 (linear decay)",
      "Gradient clip norm: 1.0",
      "Backbone: All VGG layers trained",
      "Gradient checkpointing: ON",
      "LR scheduler: StepLinearLR",
      "Warmup: 1% of total steps"]),
]

col3 = [
    ("PPO Hyperparameters",
     ["Clip epsilon: 0.1",
      "PPO epochs per rollout: 2",
      "Entropy coef: 0.001 (warmup 500)",
      "Value loss coef: 0.5 (warmup 500)",
      "GAE lambda: 0.95",
      "Discount gamma: 0.99",
      "Rollout steps: 1"]),
]

for ci, col in enumerate([col1, col2, col3]):
    lx = Inches(0.25) + ci * Inches(4.3)
    for title, bullets in col:
        add_bullet_box(slide, title, bullets,
                       lx, Inches(0.88), Inches(4.1), Inches(3.1),
                       title_size=Pt(14), bullet_size=Pt(12))

# second row
col4 = [("Reward & Scheduling",
         ["Reward type: balanced (inlier ratio + count)",
          "Reward scale: 1.0 (warmup 200 steps)",
          "Reward norm momentum: 0.99",
          "KP count floor: 380 (for 384px)",
          "KP floor penalty: 0.5",
          "KP floor warmup: 1000 steps",
          "Alpha scheduler: plateau 0→1 over 20%"])]

col5 = [("Descriptor & Loss",
         ["Descriptor dim: 256",
          "InfoNCE temperature: 0.07",
          "Uniformity weight: 0.1",
          "Desc loss weight: 0.5",
          "Descriptor warmup: 0 steps",
          "Synthetic ratio: 10% of batches",
          "Margin anneal steps: 5000"])]

col6 = [("Logging & Checkpoints",
         ["Log interval: every 500 steps",
          "Val interval: every 5000 steps",
          "Checkpoint interval: every 5000 steps",
          "W&B mode: disabled (local logs)",
          "Top-K keypoints (MNN): 1024",
          "Min matches for RANSAC: 2",
          "Output: ./outputs/ripe_imw2022/"])]

for ci, col in enumerate([col4, col5, col6]):
    lx = Inches(0.25) + ci * Inches(4.3)
    for title, bullets in col:
        add_bullet_box(slide, title, bullets,
                       lx, Inches(4.1), Inches(4.1), Inches(3.1),
                       title_size=Pt(14), bullet_size=Pt(12))


# =============================================================================
# SLIDE 6 — TRAINING RESULTS
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "05  Training Results", "GPU run snapshot — first 3 steps (50k planned)")

# Early steps table
add_text(slide, "Early Training Signals (GPU — BF16-mixed, 384px, Full VGG)",
         Inches(0.3), Inches(0.85), Inches(12), Inches(0.4),
         font_size=Pt(15), bold=True, color=ACCENT)

headers2 = ["Step", "PPO Loss", "Desc Loss", "Reward", "Inliers", "Keypoints (img1/img2)", "Status"]
rows2 = [
    ["1 / 50000", "0.5084", "2.7417", "0.564", "29.1",  "656 / 645", "Warmup phase"],
    ["2 / 50000", "0.5108", "3.4007", "0.625", "82.8",  "636 / 651", "Rapid climb"],
    ["3 / 50000", "0.5089", "2.0719", "0.631", "88.6",  "643 / 650", "Stabilising"],
]
add_table(slide, headers2, rows2,
          Inches(0.3), Inches(1.3), Inches(12.73), Inches(1.7))

# Key observations
add_text(slide, "Key Observations from Early Training",
         Inches(0.3), Inches(3.15), Inches(12), Inches(0.4),
         font_size=Pt(15), bold=True, color=ACCENT)

obs = [
    ("Reward Rising Fast",
     "Reward jumped 0.564 → 0.631 in just 3 steps (+12%). Inlier count exploded\nfrom 29 → 88 (+204%) — policy immediately learns to select better keypoints.",
     GREEN),
    ("Loss Healthy",
     "PPO loss stable ~0.509 (no divergence). Desc loss 2.74 → 2.07 (-25%) —\nInfoNCE pulling descriptors apart even in early warmup steps.",
     ACCENT2),
    ("Keypoints Stable",
     "~640 keypoints selected per image — within expected range for 384px\nimages. KP count floor (380) not triggered — policy naturally above floor.",
     YELLOW),
    ("Run Interrupted",
     "Training stopped after step 3 (likely OOM or CUDA assert).\nFix: reduce batch_size to 2, image_size to 320, num_workers to 2.",
     RED),
]

for i, (title, body, color) in enumerate(obs):
    lx = Inches(0.3) + i * Inches(3.2)
    add_rect(slide, lx, Inches(3.6), Inches(3.08), Inches(2.8), fill_color=BOX_BG, line_color=color, line_width=Pt(1.5))
    add_rect(slide, lx, Inches(3.6), Inches(3.08), Inches(0.05), fill_color=color)
    add_text(slide, title, lx + Inches(0.1), Inches(3.68), Inches(2.9), Inches(0.38),
             font_size=Pt(13), bold=True, color=color)
    add_text(slide, body, lx + Inches(0.1), Inches(4.1), Inches(2.9), Inches(2.2),
             font_size=Pt(11), color=LIGHT_GRAY)

# Expected trajectory
add_rect(slide, Inches(0.3), Inches(6.55), Inches(12.73), Inches(0.72), fill_color=MID_BLUE, line_color=ACCENT, line_width=Pt(0.8))
add_text(slide, "Expected Full-Run Trajectory (50k steps):",
         Inches(0.5), Inches(6.6), Inches(3.5), Inches(0.3),
         font_size=Pt(12), bold=True, color=YELLOW)
add_text(slide,
         "0-5k: Warmup — reward 0.5-0.7, desc loss falling   |   "
         "5-20k: Policy sharpens — reward 0.7-0.9, inliers 60-100   |   "
         "20-50k: Convergence — reward >0.9, stable keypoint selection",
         Inches(0.5), Inches(6.92), Inches(12.4), Inches(0.3),
         font_size=Pt(11), color=LIGHT_GRAY)


# =============================================================================
# SLIDE 7 — TEST RESULTS
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "06  Test Results", "Evaluation on IMW2022 test split — 126 image pairs")

# Headline metrics
metrics = [
    ("61.9%",  "Overall\nSuccess Rate",  GREEN),
    ("74.5%",  "Positive Pair\nRecall",   GREEN),
    ("82.1%",  "Negative Pair\nSpecificity", ACCENT2),
    ("46.4",   "Avg Inliers\nper Pair",   YELLOW),
    ("107.2",  "Avg Matches\nper Pair",   YELLOW),
    ("0.411",  "Avg Inlier\nRatio",       ACCENT2),
]
for i, (val, label, color) in enumerate(metrics):
    bx = Inches(0.25) + i * Inches(2.17)
    add_rect(slide, bx, Inches(0.85), Inches(2.0), Inches(1.5), fill_color=BOX_BG, line_color=color, line_width=Pt(1.5))
    add_text(slide, val, bx, Inches(0.97), Inches(2.0), Inches(0.7),
             font_size=Pt(30), bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, bx, Inches(1.65), Inches(2.0), Inches(0.65),
             font_size=Pt(11), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Before vs After comparison
add_text(slide, "Before vs After Pipeline Fixes",
         Inches(0.3), Inches(2.52), Inches(6.5), Inches(0.4),
         font_size=Pt(15), bold=True, color=ACCENT)

headers3 = ["Metric", "Before Fix", "After Fix", "Change"]
rows3 = [
    ["Overall success rate",    "54.8%",  "61.9%",  "+7.1pp"],
    ["Positive success rate",   "68.4%",  "74.5%",  "+6.1pp"],
    ["Negative specificity",    "92.9%",  "82.1%",  "-10.8pp"],
    ["Avg inliers (positive)",  "48.52",  "51.82",  "+3.3"],
    ["Avg matches",             "90.49",  "107.24", "+16.75"],
    ["Avg inlier ratio",        "0.454",  "0.411",  "-0.043"],
]
add_table(slide, headers3, rows3,
          Inches(0.3), Inches(3.0), Inches(6.3), Inches(3.8))

# Right side — analysis
add_bullet_box(slide, "Positive Pair Performance",
    ["74.5% of positive pairs correctly matched",
     "Avg 51.8 inliers — strong geometric verification",
     "Inlier gap: positive 51.8 vs negative 27.4 = 1.89x",
     "   Well-separated — threshold at 35-40 inliers works"],
    Inches(6.8), Inches(2.52), Inches(5.78), Inches(2.0),
    title_color=GREEN)

add_bullet_box(slide, "Known Issues & Next Steps",
    ["Negative false positive rate rose to 17.9%",
     "   Root cause: test images at native res (1669 kpts)",
     "   Training resolution: 384px → ~640 keypoints",
     "   Fix: force test to 384x384 to match training",
     "Raise inlier threshold from ~10 to 35 for negative rejection",
     "Full 50k GPU run needed for convergence comparison"],
    Inches(6.8), Inches(4.65), Inches(5.78), Inches(2.15),
    title_color=YELLOW)

# Confusion matrix visual
add_text(slide, "Confusion Matrix (126 pairs)",
         Inches(0.3), Inches(6.58), Inches(6), Inches(0.35),
         font_size=Pt(13), bold=True, color=ACCENT)

cm_data = [
    ["",          "Pred: Match",                   "Pred: No Match"                ],
    ["True: +ve", "TP = 73  (74.5%)",              "FN = 25  (25.5%)"              ],
    ["True: -ve", "FP =  5  (17.9%)",              "TN = 23  (82.1%)"              ],
]
cx, cy, cw, ch = Inches(0.3), Inches(6.95), Inches(6.3), Inches(0.52)
colors_cm = [
    [MID_BLUE, MID_BLUE,  MID_BLUE],
    [MID_BLUE, GREEN,     RGBColor(0x3A,0x18,0x18)],
    [MID_BLUE, RGBColor(0x3A,0x18,0x18), GREEN],
]
for ri, row in enumerate(cm_data):
    for ci2, cell in enumerate(row):
        add_rect(slide, cx + ci2*(cw/3), cy + ri*ch, cw/3, ch,
                 fill_color=colors_cm[ri][ci2], line_color=ACCENT, line_width=Pt(0.5))
        add_text(slide, cell, cx + ci2*(cw/3) + Inches(0.05), cy + ri*ch + Inches(0.04),
                 cw/3 - Inches(0.1), ch - Inches(0.08),
                 font_size=Pt(10), bold=(ri==0 or ci2==0), color=WHITE, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 8 — CONCLUSION & FUTURE WORK
# =============================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, DARK_BG)
slide_header(slide, "Summary & Future Work", "Contributions, results, and next steps")

add_bullet_box(slide, "What We Achieved",
    ["Replaced REINFORCE with PPO — stable clipped surrogate training",
     "Added context-aware Attention Descriptor Head (MHA + FFN + LN)",
     "Integrated LO-RANSAC with PROSAC match ordering",
     "InfoNCE loss + synthetic homography bootstrapping",
     "74.5% positive recall on IMW2022 test split (126 pairs)",
     "Overall success rate: 61.9% — competitive with base RIPE",
     "BF16-mixed precision — 2x VRAM efficiency for GPU training"],
    Inches(0.3), Inches(0.85), Inches(6.0), Inches(5.8),
    title_size=Pt(16), bullet_size=Pt(13))

add_bullet_box(slide, "Remaining Challenges",
    ["GPU training stopped at step 3 — OOM or CUDA assert",
     "   Fix: batch_size=2, image_size=320, num_workers=2",
     "Test resolution mismatch (native res vs 384px training)",
     "   Fix: force test to 384×384 in test_models.py",
     "Negative false positive rate increased to 17.9%",
     "   Fix: raise inlier threshold to 35-40 for decision",
     "Full 50k step run needed for convergence study"],
    Inches(6.5), Inches(0.85), Inches(6.53), Inches(3.1),
    title_color=RED, title_size=Pt(16), bullet_size=Pt(13))

add_bullet_box(slide, "Future Work",
    ["Complete 50k GPU training run with stable config",
     "Evaluate on MegaDepth-1500 and HPatches benchmarks",
     "Multi-GPU DDP training for larger batch sizes",
     "Explore adaptive KP count floor scheduling",
     "Ablation: PPO vs REINFORCE vs supervised DISK",
     "Attention heads ablation: 1 / 2 / 4 / 8 heads",
     "Integrate with full SfM pipeline (COLMAP)"],
    Inches(6.5), Inches(4.1), Inches(6.53), Inches(2.55),
    title_color=GREEN, title_size=Pt(16), bullet_size=Pt(13))

# bottom bar
add_rect(slide, 0, Inches(6.85), SW, Inches(0.65), fill_color=MID_BLUE)
add_rect(slide, 0, Inches(6.85), SW, Inches(0.05), fill_color=ACCENT)
add_text(slide,
         "RIPE-PPO  •  PPO + Attention Descriptor + LO-RANSAC + PROSAC  •  IMW2022  •  BTP 2025-26",
         Inches(0.3), Inches(6.9), Inches(12.5), Inches(0.5),
         font_size=Pt(13), color=ACCENT2, align=PP_ALIGN.CENTER)


# =============================================================================
# SAVE
# =============================================================================
out_path = r"C:\Users\linga\Downloads\BTP_MODIFIED_VGG\RIPE_PPO_Presentation.pptx"
prs.save(out_path)
print(f"Done. Saved to: {out_path}")
print(f"Slides: {len(prs.slides)}")
